import importlib
import os
import unittest
from unittest.mock import AsyncMock, patch

from pptx import Presentation

from reactor_tool.model.protocal import ReportRequest


tool_api = importlib.import_module("reactor_tool.api.tool")


class ReportPptxTest(unittest.IsolatedAsyncioTestCase):
    async def test_non_stream_report_uploads_real_pptx(self):
        response, upload = await self._run_report(stream=False)

        self.assertEqual("demo.pptx", response["fileInfo"][0]["fileName"])
        upload.assert_awaited_once()

    async def test_stream_report_uploads_real_pptx(self):
        _, upload = await self._run_report(stream=True)

        upload.assert_awaited_once()

    async def test_report_honors_explicit_slide_limit_and_keeps_final_section(self):
        async def fake_report(**kwargs):
            yield "\n---\n".join(
                f"## Slide {index}\n\n- Point {index}" for index in range(1, 9)
            )

        async def inspect_upload(file_path, _request_id):
            presentation = Presentation(file_path)
            self.assertEqual(3, len(presentation.slides))
            titles = [slide.shapes[1].text for slide in presentation.slides]
            self.assertEqual(["Slide 1", "Slide 2", "Slide 8"], titles)
            return {"fileName": "limited.pptx", "ossUrl": "https://files.test/limited.pptx"}

        body = ReportRequest.model_validate({
            "requestId": "ppt-limited",
            "query": "生成 3 页以内的演示文稿",
            "task": "Only use these sections. 第 1 页为封面，第 2 页为正文，第 3 页为结论。",
            "fileName": "limited",
            "fileType": "ppt",
        })
        with patch("reactor_tool.tool.report.report", side_effect=fake_report), \
                patch.object(tool_api, "upload_file_by_path", new=AsyncMock(side_effect=inspect_upload)):
            await tool_api.post_report(body)

    async def test_report_ignores_preface_and_keeps_final_boundary_points(self):
        async def fake_report(**kwargs):
            slides = [f"## 第 {index} 页\n\n- 内容 {index}" for index in range(1, 6)]
            slides.append(
                "## 第 6 页｜边界声明\n\n"
                "### 证据范围\n"
                "- 当前证据为本地验证\n"
                "- 不是生产 SLA\n"
                "- 离线确定性组件回归\n"
                "- 在线真实模型流程\n"
                "- 会话 owner 隔离\n"
                "- 上下文按预算压缩\n"
                "- DNS rebinding 仍是边界\n"
                "- ASK 人工审批通路仍是边界"
            )
            yield "无法直接附加二进制文件。\n\n---\n\n" + "\n\n---\n\n".join(slides)

        async def inspect_upload(file_path, _request_id):
            presentation = Presentation(file_path)
            self.assertEqual(6, len(presentation.slides))
            titles = [slide.shapes[1].text for slide in presentation.slides]
            self.assertEqual("第 1 页", titles[0])
            self.assertEqual("第 6 页｜边界声明", titles[-1])
            last_slide = "\n".join(
                shape.text for shape in presentation.slides[-1].shapes if hasattr(shape, "text")
            )
            for required in ("本地验证", "生产 SLA", "DNS rebinding", "ASK"):
                self.assertIn(required, last_slide)
            return {"fileName": "bounded.pptx", "ossUrl": "https://files.test/bounded.pptx"}

        body = ReportRequest.model_validate({
            "requestId": "ppt-boundary",
            "query": "生成 6 页以内的演示文稿",
            "task": "最后一页必须保留全部边界声明",
            "fileName": "bounded",
            "fileType": "ppt",
        })
        with patch("reactor_tool.tool.report.report", side_effect=fake_report), \
                patch.object(tool_api, "upload_file_by_path", new=AsyncMock(side_effect=inspect_upload)):
            await tool_api.post_report(body)

    async def _run_report(self, stream):
        async def fake_report(**kwargs):
            yield "## First slide\n\n- Real PPTX output\n\n---\n\n## Second slide\n\n- Downloadable artifact"

        async def inspect_upload(file_path, request_id):
            self.assertEqual("ppt-request", request_id)
            self.assertTrue(os.path.isfile(file_path))
            presentation = Presentation(file_path)
            self.assertEqual(2, len(presentation.slides))
            return {
                "fileName": os.path.basename(file_path),
                "ossUrl": "https://files.test/demo.pptx",
                "domainUrl": "https://files.test/preview/demo.pptx",
                "fileSize": os.path.getsize(file_path),
            }

        body = ReportRequest.model_validate({
            "requestId": "ppt-request",
            "query": "Create a presentation",
            "task": "Create two slides",
            "fileName": "demo",
            "fileType": "ppt",
            "stream": stream,
        })
        upload = AsyncMock(side_effect=inspect_upload)
        with patch("reactor_tool.tool.report.report", side_effect=fake_report), \
                patch.object(tool_api, "upload_file_by_path", new=upload):
            response = await tool_api.post_report(body)
            if stream:
                async for _ in response.body_iterator:
                    pass
        return response, upload


if __name__ == "__main__":
    unittest.main()
