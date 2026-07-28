# -*- coding: utf-8 -*-
import asyncio
import io
import os
import tempfile
import unittest
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock, patch

from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from starlette.datastructures import UploadFile

from reactor_tool.api.file_manage import router
from reactor_tool.db.file_table_op import FileDB, FileInfoOp


class FileManageApiTest(unittest.TestCase):
    def setUp(self):
        app = FastAPI()
        app.include_router(router, prefix="/v1/file_tool")
        self.client = TestClient(app)

    def test_should_preview_file_when_url_contains_nested_path_segments(self):
        with tempfile.NamedTemporaryFile("w", encoding="utf-8", suffix=".md", delete=False) as temp_file:
            temp_file.write("# 测试文件\n")
            file_path = temp_file.name

        try:
            file_info = SimpleNamespace(file_path=file_path)
            with patch(
                "reactor_tool.api.file_manage.FileInfoOp.get_by_file_id",
                new=AsyncMock(side_effect=[file_info]),
            ) as get_by_file_id:
                response = self.client.get(
                    "/v1/file_tool/preview/session-001/colbymchenry/demo.md"
                )

            self.assertEqual(200, response.status_code)
            self.assertEqual("# 测试文件", response.text.strip())
            get_by_file_id.assert_awaited_once()
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)

    def test_should_fallback_to_legacy_file_id_for_nested_path_segments(self):
        with tempfile.NamedTemporaryFile("w", encoding="utf-8", suffix=".md", delete=False) as temp_file:
            temp_file.write("# 历史文件\n")
            file_path = temp_file.name

        try:
            file_info = SimpleNamespace(file_path=file_path)
            with patch(
                "reactor_tool.api.file_manage.FileInfoOp.get_by_file_id",
                new=AsyncMock(side_effect=[None, file_info]),
            ) as get_by_file_id:
                response = self.client.get(
                    "/v1/file_tool/preview/session-002/colbymchenry/legacy.md"
                )

            self.assertEqual(200, response.status_code)
            self.assertEqual("# 历史文件", response.text.strip())
            self.assertEqual(2, get_by_file_id.await_count)
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)

    def test_should_render_pptx_as_safe_html_preview(self):
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_file:
            file_path = temp_file.name
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(0.8)).text = "demo.pptx"
        body = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(10), Inches(5)).text_frame
        body.paragraphs[0].text = "Slide 1｜项目技术报告"
        for index in range(1, 17):
            body.add_paragraph().text = f"Agent Loop 证据 {index}"
        presentation.save(file_path)

        try:
            file_info = SimpleNamespace(file_path=file_path)
            with patch(
                "reactor_tool.api.file_manage.FileInfoOp.get_by_file_id",
                new=AsyncMock(return_value=file_info),
            ):
                response = self.client.get("/v1/file_tool/preview/session-003/demo.pptx")

            self.assertEqual(200, response.status_code)
            self.assertTrue(response.headers["content-type"].startswith("text/html"))
            self.assertIn("<h1>Slide 1｜项目技术报告</h1>", response.text)
            self.assertIn('class="slide very-dense"', response.text)
            self.assertNotIn("<script", response.text.lower())
        finally:
            if os.path.exists(file_path):
                os.remove(file_path)

    def test_should_store_binary_upload_under_session_directory(self):
        with tempfile.TemporaryDirectory(prefix="file-manage-local-") as temp_dir:
            original_work_dir = FileDB._work_dir
            FileDB._work_dir = temp_dir
            upload_file = UploadFile(filename="poster.png", file=io.BytesIO(b"fake-image-bytes"))

            try:
                with patch.object(
                    FileInfoOp,
                    "add",
                    new=AsyncMock(side_effect=lambda file_info: file_info),
                ):
                    file_info = asyncio.run(
                        FileInfoOp.add_by_file(
                            file=upload_file,
                            file_id="file-id-001",
                            request_id="session-1779798194080-9667",
                        )
                    )

                saved_path = Path(file_info.file_path)
                self.assertTrue(saved_path.exists())
                self.assertEqual("session-1779798194080-9667", saved_path.parent.name)
                self.assertEqual("poster.png", saved_path.name)
            finally:
                FileDB._work_dir = original_work_dir

    def test_should_delete_object_before_metadata(self):
        with patch.object(FileInfoOp, "delete", new=AsyncMock()) as delete:
            response = self.client.delete("/v1/file_tool/file-id-001")

        self.assertEqual(200, response.status_code)
        self.assertEqual({"fileId": "file-id-001", "deleted": True}, response.json())
        delete.assert_awaited_once_with("file-id-001")
