from html import escape
import re
import textwrap
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


def render_pptx(content: str, output_file: str, max_slides: int | None = None) -> str:
    """Render report HTML or Markdown into a real, portable PPTX file."""
    slides = _extract_slides(content, max_slides)
    if not slides:
        raise ValueError("PPT content does not contain any usable slides")

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    for page_number, (title, points) in enumerate(slides, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(17, 24, 32)

        accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(45, 212, 191)
        accent.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = title
        title_paragraph.font.name = "Microsoft YaHei"
        title_paragraph.font.size = Pt(28)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(245, 247, 250)

        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.45), Inches(4.95))
        body_frame = body_box.text_frame
        body_frame.clear()
        body_frame.word_wrap = True
        body_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        body_frame.margin_left = Inches(0.08)
        body_frame.margin_right = Inches(0.08)
        body_frame.margin_top = Inches(0.05)
        body_frame.margin_bottom = Inches(0.05)
        for index, point in enumerate(points or [" "]):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = f"- {point}" if point.strip() else " "
            paragraph.font.name = "Microsoft YaHei"
            paragraph.font.size = Pt(14 if len(points) > 10 else 17 if len(points) > 7 else 19 if len(points) > 5 else 21)
            paragraph.font.color.rgb = RGBColor(220, 226, 232)
            paragraph.space_after = Pt(6 if len(points) > 7 else 12)

        footer = slide.shapes.add_textbox(Inches(11.6), Inches(6.85), Inches(0.8), Inches(0.3))
        footer_paragraph = footer.text_frame.paragraphs[0]
        footer_paragraph.text = str(page_number)
        footer_paragraph.alignment = PP_ALIGN.RIGHT
        footer_paragraph.font.size = Pt(11)
        footer_paragraph.font.color.rgb = RGBColor(148, 163, 184)

    target = Path(output_file)
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(target)
    return str(target)


def safe_pptx_name(file_name: str) -> str:
    name = Path((file_name or "report").strip()).name
    if name.lower().endswith(".pptx"):
        name = name[:-5]
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name).strip(" .") or "report"
    return f"{name}.pptx"


def render_pptx_preview_html(file_path: str) -> str:
    """Render the text slides in a generated PPTX as a safe browser preview."""
    presentation = Presentation(file_path)
    if not presentation.slides:
        raise ValueError("PPTX does not contain any slides")

    pages = []
    for page_number, slide in enumerate(presentation.slides, start=1):
        blocks = []
        for shape in sorted(slide.shapes, key=lambda item: (item.top, item.left)):
            if not getattr(shape, "has_text_frame", False):
                continue
            lines = [
                re.sub(r"^\s*[-•]\s*", "", paragraph.text).strip()
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text.strip()
            ]
            if lines and lines != [str(page_number)]:
                blocks.append(lines)

        title = blocks[0][0] if blocks else f"第 {page_number} 页"
        points = (blocks[0][1:] if blocks else []) + [line for block in blocks[1:] for line in block]
        if title.lower().endswith(".pptx") and points and re.match(r"(?i)^slide\s*\d+", points[0]):
            title = points.pop(0)
        body = "".join(f"<li>{escape(point)}</li>" for point in points)
        if not body:
            body = '<li class="empty">此页没有可提取的文本内容</li>'
        density = " very-dense" if len(points) > 14 else " dense" if len(points) > 10 else ""
        pages.append(
            f'<section class="slide-wrap" aria-label="第 {page_number} 页">'
            f'<article class="slide{density}">'
            '<span class="accent" aria-hidden="true"></span>'
            f'<h1>{escape(title)}</h1><ul>{body}</ul>'
            f'<span class="page-number">{page_number}</span>'
            '</article></section>'
        )

    return """<!doctype html>
<html lang="zh-CN">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>PPTX Preview</title>
<style>
* { box-sizing: border-box; }
html { background: #eef0f2; color-scheme: light; }
body { margin: 0; font-family: "Microsoft YaHei", "Segoe UI", sans-serif; background: #eef0f2; color: #f5f7fa; }
.deck { display: grid; gap: 24px; padding: 24px; }
.slide-wrap { display: flex; justify-content: center; }
.slide { position: relative; width: min(100%, 1100px); aspect-ratio: 16 / 9; overflow: hidden; border: 1px solid #d8dde3; background: #111820; box-shadow: 0 10px 28px rgba(17, 24, 32, .14); }
.accent { position: absolute; inset: 0 auto 0 0; width: 1.2%; background: #2dd4bf; }
h1 { margin: 4.2% 7% 0; font-size: 30px; line-height: 1.25; font-weight: 700; }
ul { display: grid; gap: 12px; margin: 5% 8% 0; padding: 0; list-style: none; }
li { position: relative; padding-left: 22px; color: #dce2e8; font-size: 18px; line-height: 1.42; }
li::before { position: absolute; left: 0; content: "•"; color: #2dd4bf; }
li.empty { color: #94a3b8; }
.slide.dense h1 { font-size: 26px; }
.slide.dense ul { gap: 6px; margin-top: 3.5%; }
.slide.dense li { font-size: 15px; line-height: 1.3; }
.slide.very-dense h1 { font-size: 24px; }
.slide.very-dense ul { gap: 4px; margin-top: 3%; }
.slide.very-dense li { font-size: 13px; line-height: 1.2; }
.page-number { position: absolute; right: 6.8%; bottom: 4.5%; color: #94a3b8; font-size: 12px; }
@media (max-width: 640px) {
  .deck { gap: 12px; padding: 12px; }
  h1 { font-size: 20px; }
  ul { gap: 6px; margin-top: 4%; }
  li { padding-left: 16px; font-size: 13px; line-height: 1.3; }
  .slide.dense h1, .slide.very-dense h1 { font-size: 18px; }
  .slide.dense ul, .slide.very-dense ul { gap: 2px; margin-top: 2%; }
  .slide.dense li { font-size: 11px; line-height: 1.18; }
  .slide.very-dense li { font-size: 10px; line-height: 1.1; }
}
</style>
</head>
<body><main class="deck">""" + "".join(pages) + "</main></body></html>"


def requested_slide_limit(*instructions: str | None) -> int | None:
    text = "\n".join(value for value in instructions if value)
    patterns = (
        r"(?:最多|不超过|不多于|上限(?:为)?|页数[:：]?\s*)(\d{1,2})\s*(?:页|张)",
        r"(\d{1,2})\s*(?:页|张)\s*(?:以内|以下|之内)",
        r"(?i)(?:no\s+more\s+than|up\s+to|max(?:imum)?(?:\s+of)?)\s+(\d{1,2})\s+slides?",
    )
    matches = [int(value) for pattern in patterns for value in re.findall(pattern, text)]
    return min((value for value in matches if value > 0), default=None)


def _extract_slides(content: str, max_slides: int | None = None) -> list[tuple[str, list[str]]]:
    normalized = (content or "").strip()
    if not normalized:
        return []
    if re.search(r"<(?:!doctype|html|body|section|div)\b", normalized, re.IGNORECASE):
        raw_slides = _extract_html_slides(normalized)
    else:
        raw_slides = _extract_markdown_slides(normalized)

    if max_slides and len(raw_slides) >= max_slides:
        if len(raw_slides) > max_slides:
            raw_slides = raw_slides[:max_slides - 1] + [raw_slides[-1]]
        return raw_slides

    slides = []
    for title, points in raw_slides:
        chunks = _chunk_points(points)
        if max_slides:
            chunks = chunks[:max(1, max_slides - len(slides))]
        for index, chunk in enumerate(chunks, start=1):
            page_title = title if index == 1 else f"{title} (continued)"
            slides.append((page_title, chunk))
            if max_slides and len(slides) >= max_slides:
                return slides
    return slides


def _extract_html_slides(content: str) -> list[tuple[str, list[str]]]:
    soup = BeautifulSoup(content, "html.parser")
    elements = soup.select(".slide") or ([soup.body] if soup.body else [soup])
    slides = []
    for index, element in enumerate(elements, start=1):
        heading = element.find(["h1", "h2", "h3"])
        title = _clean_text(heading.get_text(" ", strip=True)) if heading else f"Slide {index}"
        points = []
        for node in element.find_all(["p", "li", "h4", "h5"]):
            if node.find_parent("li") is not None and node.name != "li":
                continue
            text = _clean_text(node.get_text(" ", strip=True))
            if text and text != title and text not in points:
                points.append(text)
        slides.append((title, points))
    return slides


def _extract_markdown_slides(content: str) -> list[tuple[str, list[str]]]:
    sections = re.split(r"(?m)^\s*---\s*$", content)
    has_headings = any(re.search(r"(?m)^\s*#{1,6}\s+", section) for section in sections)
    slides = []
    for index, section in enumerate(sections, start=1):
        lines = [line.strip() for line in section.splitlines() if line.strip()]
        if not lines:
            continue
        heading_index = next((i for i, line in enumerate(lines) if re.match(r"^#{1,6}\s+", line)), None)
        if heading_index is None and has_headings:
            continue
        title = _clean_text(re.sub(r"^#{1,6}\s+", "", lines[heading_index])) if heading_index is not None else f"Slide {index}"
        points = []
        for i, line in enumerate(lines):
            if i == heading_index or line.startswith("```"):
                continue
            text = _clean_text(re.sub(r"^(?:#{1,6}\s+|[-*+]\s+|\d+[.)]\s+)", "", line))
            if text and text not in points:
                points.append(text)
        slides.append((title, points))
    return slides


def _chunk_points(points: list[str]) -> list[list[str]]:
    expanded = []
    for point in points:
        expanded.extend(part.strip() for part in textwrap.wrap(point, width=100) if part.strip())
    if not expanded:
        return [[]]

    chunks = []
    current = []
    current_chars = 0
    for point in expanded:
        if current and (len(current) >= 7 or current_chars + len(point) > 560):
            chunks.append(current)
            current = []
            current_chars = 0
        current.append(point)
        current_chars += len(point)
    if current:
        chunks.append(current)
    return chunks


def _clean_text(text: str) -> str:
    normalized = re.sub(r"[`*_~]", "", text or "")
    return re.sub(r"\s+", " ", normalized).strip()
